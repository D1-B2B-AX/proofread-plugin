"""
PPT 교안 텍스트 추출 스크립트
- .pptx 파일에서 슬라이드별 텍스트를 추출
- 텍스트 박스, 표, 그룹 도형 내 텍스트 모두 포함
- 각 텍스트의 슬라이드 내 위치(상/중/하) 정보 포함
- 반복 텍스트 변형 감지 (유사 문장 중 소수 변형 자동 플래그)
- 형태소 분석 기반 의심 단어 감지 (사전에 없는 한글 단어 자동 플래그)
- 스피커 노트 제외
"""

import sys
import os
import json
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 공통 감지 유틸리티 import
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from detection_utils import detect_text_variants, detect_suspicious_words, detect_suspicious_english, detect_broken_text

# 슬라이드 번호 자리표시자 패턴 (‹#›, <#>, <<#>> 등)
SLIDE_NUMBER_PATTERN = re.compile(r'^[‹<«\[]*#[›>\»\]]*$')


def get_position_label(top, total_height):
    """도형의 top 좌표와 슬라이드 전체 높이로 상/중/하 판별"""
    ratio = top / total_height if total_height > 0 else 0
    if ratio < 0.33:
        return "상"
    elif ratio < 0.66:
        return "중"
    else:
        return "하"


def extract_text_from_shape(shape, slide_height):
    """도형에서 텍스트 추출 (재귀적으로 그룹 도형 내부도 탐색)"""
    items = []
    shape_top = shape.top if shape.top is not None else 0
    position = get_position_label(shape_top, slide_height)

    # 일반 텍스트 박스 / 도형
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            # 슬라이드 번호 자리표시자 제외
            if text and not SLIDE_NUMBER_PATTERN.match(text):
                items.append({"text": text, "position": position})

    # 표 (Table)
    if shape.has_table:
        for row in shape.table.rows:
            row_texts = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_texts.append(cell_text)
            if row_texts:
                items.append({"text": " | ".join(row_texts), "position": position})

    # 그룹 도형 (내부 도형을 재귀 탐색)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            child_items = extract_text_from_shape(child_shape, slide_height)
            items.extend(child_items)

    return items


def extract_pptx(file_path):
    """
    PPTX 파일에서 슬라이드별 텍스트를 추출

    Returns:
        dict: {
            "file_name": 파일명,
            "total_slides": 총 슬라이드 수,
            "slides": [
                {
                    "slide_number": 슬라이드 번호,
                    "texts": [
                        {"text": 텍스트, "position": "상/중/하"},
                        ...
                    ]
                },
                ...
            ]
        }
    """
    prs = Presentation(file_path)
    slide_height = prs.slide_height  # 슬라이드 전체 높이 (EMU 단위, 가로/세로 자동 대응)

    result = {
        "file_name": file_path.split("/")[-1].split("\\")[-1],
        "total_slides": len(prs.slides),
        "slides": []
    }

    for i, slide in enumerate(prs.slides, 1):
        slide_items = []

        for shape in slide.shapes:
            items = extract_text_from_shape(shape, slide_height)
            slide_items.extend(items)

        result["slides"].append({
            "slide_number": i,
            "texts": slide_items
        })

    # 반복 텍스트 변형 감지 (줄 단위)
    result["suspected_variants"] = detect_text_variants(result["slides"], key="slide_number")

    # 형태소 분석 기반 의심 단어 감지
    result["suspicious_words"] = detect_suspicious_words(result["slides"], key="slide_number")

    # 영문 스펠 체크 기반 의심 단어 감지
    result["suspicious_english"] = detect_suspicious_english(result["slides"], key="slide_number")

    # 텍스트 깨짐 감지
    result["broken_texts"] = detect_broken_text(result["slides"], key="slide_number")

    return result




def format_output(result):
    """추출 결과를 읽기 쉬운 텍스트 형식으로 변환"""
    output_lines = []
    output_lines.append(f"파일명: {result['file_name']}")
    output_lines.append(f"총 슬라이드: {result['total_slides']}p")
    output_lines.append("")

    for slide in result["slides"]:
        slide_num = slide["slide_number"]
        texts = slide["texts"]

        output_lines.append(f"=== 슬라이드 {slide_num} ===")
        if texts:
            for item in texts:
                output_lines.append(f"[{item['position']}] {item['text']}")
        else:
            output_lines.append("(텍스트 없음)")
        output_lines.append("")

    return "\n".join(output_lines)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("사용법: python extract_pptx.py <파일경로.pptx> [--json]")
        sys.exit(1)

    file_path = sys.argv[1]
    use_json = "--json" in sys.argv

    result = extract_pptx(file_path)

    if use_json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        print(format_output(result))
